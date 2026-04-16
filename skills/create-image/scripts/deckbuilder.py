#!/usr/bin/env python3
"""Banana Claude -- Slide Deck Builder

Assembles generated slide images into an editable .pptx presentation with
text layers, logos, and brand styling.

Usage:
    deckbuilder.py build --images DIR --output deck.pptx [--preset NAME] [--title "Title"] [--layout standard|fullbleed|split] [--logo PATH]
    deckbuilder.py estimate --images DIR
    deckbuilder.py layouts
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

# Graceful import of python-pptx (already approved in this project)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print(json.dumps({
        "error": True,
        "message": "python-pptx required. Install: pip install python-pptx",
    }))
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None  # non-critical -- used only for logo aspect ratio

# Sibling import -- pantone_lookup lives next to this script
_SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(_SCRIPT_DIR))
from pantone_lookup import hex_to_rgb

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

PRESETS_DIR = Path.home() / ".banana" / "presets"

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}

SLIDE_WIDTH = Inches(13.333)   # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

LAYOUTS = {
    "fullbleed": "Image fills entire slide, text overlaid with semi-transparent backdrop",
    "standard": "Image occupies top 60%, text area bottom 40% with brand background",
    "split": "50/50 — image on left, text content on right",
}

DEFAULT_LAYOUT = "standard"

# Font mapping from typography descriptions to safe built-in fonts
_FONT_MAP = {
    "bold":         "Arial Black",
    "condensed":    "Arial Black",
    "modern":       "Calibri",
    "serif":        "Georgia",
    "classic":      "Georgia",
    "elegant":      "Georgia",
    "monospace":    "Courier New",
    "clean":        "Calibri",
    "sans":         "Arial",
    "sans-serif":   "Arial",
}

DEFAULT_HEADING_FONT = "Arial Black"
DEFAULT_BODY_FONT = "Calibri"

# Default brand colors (dark theme)
DEFAULT_BG_COLOR = "#1A1A2E"
DEFAULT_ACCENT_COLOR = "#E94560"
DEFAULT_TEXT_COLOR = "#FFFFFF"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _hex_to_pptx_rgb(hex_color: str) -> RGBColor:
    """Convert '#FF0000' or 'FF0000' to RGBColor(255, 0, 0)."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return RGBColor(255, 255, 255)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _map_font(typography_desc: str) -> tuple:
    """Map a preset typography description to (heading_font, body_font)."""
    if not typography_desc:
        return DEFAULT_HEADING_FONT, DEFAULT_BODY_FONT
    desc = typography_desc.lower()
    heading = DEFAULT_HEADING_FONT
    body = DEFAULT_BODY_FONT
    for keyword, font in _FONT_MAP.items():
        if keyword in desc:
            heading = font
            break
    # Body stays clean/readable regardless
    if "serif" in desc and "sans" not in desc:
        body = "Georgia"
    return heading, body


def _load_slide_images(images_dir: Path) -> list:
    """Scan directory for slide images, sorted by filename.

    Returns list of dicts: {number, name, path}
    Expects filenames like slide-01-title.png or just 01.png.
    """
    images = []
    for f in sorted(images_dir.iterdir()):
        if f.suffix.lower() in IMAGE_EXTENSIONS and f.is_file():
            stem = f.stem
            # Try to extract number and name from patterns like "slide-01-title"
            match = re.match(r"(?:slide-?)?(\d+)[-_]?(.*)", stem, re.IGNORECASE)
            if match:
                number = int(match.group(1))
                name = match.group(2).replace("-", " ").replace("_", " ").strip()
            else:
                number = len(images) + 1
                name = stem.replace("-", " ").replace("_", " ").strip()
            images.append({"number": number, "name": name or f"Slide {number}", "path": str(f)})
    return images


def _load_summary(images_dir: Path) -> dict | None:
    """Read generation-summary.json if present (output by slides.py)."""
    summary_path = images_dir / "generation-summary.json"
    if summary_path.exists():
        try:
            with open(summary_path, "r") as fh:
                return json.load(fh)
        except (json.JSONDecodeError, OSError):
            pass
    return None


def _load_preset(name: str) -> dict | None:
    """Load brand preset JSON from ~/.banana/presets/."""
    if not name:
        return None
    preset_path = PRESETS_DIR / f"{name}.json"
    if not preset_path.exists():
        print(f"Warning: preset '{name}' not found at {preset_path}", file=sys.stderr)
        return None
    try:
        with open(preset_path, "r") as fh:
            return json.load(fh)
    except (json.JSONDecodeError, OSError) as exc:
        print(f"Warning: failed to load preset '{name}': {exc}", file=sys.stderr)
        return None


def _get_preset_colors(preset: dict | None) -> tuple:
    """Extract (bg_color, accent_color, text_color) from preset."""
    if not preset:
        return DEFAULT_BG_COLOR, DEFAULT_ACCENT_COLOR, DEFAULT_TEXT_COLOR
    colors = preset.get("colors", [])
    bg = colors[0] if len(colors) > 0 else DEFAULT_BG_COLOR
    accent = colors[1] if len(colors) > 1 else DEFAULT_ACCENT_COLOR
    return bg, accent, DEFAULT_TEXT_COLOR


def _get_summary_for_slide(summary_data: dict | None, slide_number: int) -> dict:
    """Get prompt/metadata for a specific slide from summary data."""
    if not summary_data:
        return {}
    slides = summary_data.get("slides", [])
    for s in slides:
        if s.get("number") == slide_number:
            return s
    return {}


def _add_logo(slide, prs, logo_path: str, preset: dict | None):
    """Add logo image to a slide."""
    if not logo_path or not Path(logo_path).exists():
        return
    logo_height = int(prs.slide_height * 0.05)  # 5% of slide height
    placement = "top-right"
    if preset and preset.get("logo_placement"):
        placement = preset["logo_placement"]

    # Determine logo dimensions preserving aspect ratio
    logo_width = logo_height  # default square assumption
    if Image:
        try:
            with Image.open(logo_path) as img:
                w, h = img.size
                if h > 0:
                    logo_width = int(logo_height * (w / h))
        except Exception:
            pass

    # Position based on placement
    margin = Inches(0.4)
    if "left" in placement:
        left = margin
    else:
        left = prs.slide_width - logo_width - margin
    if "bottom" in placement:
        top = prs.slide_height - logo_height - margin
    else:
        top = margin

    slide.shapes.add_picture(logo_path, left, top, logo_width, logo_height)


def _set_slide_bg(slide, hex_color: str):
    """Set solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_pptx_rgb(hex_color)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_title_slide(prs, title: str, subtitle: str, preset: dict | None, logo_path: str):
    """Create the opening title slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color, accent_color, text_color = _get_preset_colors(preset)
    _set_slide_bg(slide, bg_color)

    heading_font, body_font = _map_font(preset.get("typography", "") if preset else "")

    # Title text
    left = Inches(1.5)
    top = Inches(2.2)
    width = prs.slide_width - Inches(3)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title or "Presentation"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = _hex_to_pptx_rgb(accent_color)
    p.font.name = heading_font
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        sub_top = Inches(4.5)
        sub_height = Inches(1.2)
        txBox2 = slide.shapes.add_textbox(left, sub_top, width, sub_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = _hex_to_pptx_rgb(text_color)
        p2.font.name = body_font
        p2.alignment = PP_ALIGN.LEFT

    _add_logo(slide, prs, logo_path, preset)
    return slide


def _add_content_slide_fullbleed(prs, slide, image_info: dict, preset, logo_path, summary_entry):
    """Full-bleed layout: image covers slide, text overlaid at bottom."""
    # Image fills entire slide
    slide.shapes.add_picture(image_info["path"], 0, 0, prs.slide_width, prs.slide_height)

    _, accent_color, text_color = _get_preset_colors(preset)

    # Semi-transparent dark rectangle at bottom for text readability
    backdrop_top = prs.slide_height - Inches(2)
    backdrop_height = Inches(2)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        0, backdrop_top, prs.slide_width, backdrop_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.fore_color.brightness = -0.5
    shape.line.fill.background()
    # Make semi-transparent via XML manipulation
    from pptx.oxml.ns import qn
    sp_pr = shape._element.spPr
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = sp_pr.find(".//" + qn("a:solidFill"))
    if solid_fill is not None:
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is not None:
            from lxml import etree
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", "60000")  # 60% opacity

    heading_font, _ = _map_font(preset.get("typography", "") if preset else "")

    # Heading text
    left = Inches(0.8)
    top = prs.slide_height - Inches(1.6)
    width = prs.slide_width - Inches(1.6)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = image_info["name"].title()
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = heading_font
    p.alignment = PP_ALIGN.LEFT

    _add_logo(slide, prs, logo_path, preset)


def _add_content_slide_standard(prs, slide, image_info: dict, preset, logo_path, summary_entry):
    """Standard layout: image top 60%, text area bottom 40%."""
    bg_color, accent_color, text_color = _get_preset_colors(preset)
    _set_slide_bg(slide, bg_color)

    # Image in top 60%
    img_height = Inches(4.5)
    slide.shapes.add_picture(image_info["path"], 0, 0, prs.slide_width, img_height)

    heading_font, body_font = _map_font(preset.get("typography", "") if preset else "")

    # Heading in bottom area
    left = Inches(0.8)
    top = Inches(4.8)
    width = prs.slide_width - Inches(1.6)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = image_info["name"].title()
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_pptx_rgb(accent_color)
    p.font.name = heading_font
    p.alignment = PP_ALIGN.LEFT

    # Notes text if available from summary
    notes_text = summary_entry.get("prompt", "")
    if notes_text:
        notes_top = Inches(5.8)
        notes_height = Inches(1.4)
        txBox2 = slide.shapes.add_textbox(left, notes_top, width, notes_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = notes_text[:200]  # truncate long prompts
        p2.font.size = Pt(16)
        p2.font.color.rgb = _hex_to_pptx_rgb(text_color)
        p2.font.name = body_font
        p2.alignment = PP_ALIGN.LEFT

    _add_logo(slide, prs, logo_path, preset)


def _add_content_slide_split(prs, slide, image_info: dict, preset, logo_path, summary_entry):
    """Split layout: image left 50%, text right 50%."""
    bg_color, accent_color, text_color = _get_preset_colors(preset)
    _set_slide_bg(slide, bg_color)

    half_width = int(prs.slide_width / 2)

    # Image on left half
    slide.shapes.add_picture(image_info["path"], 0, 0, half_width, prs.slide_height)

    heading_font, body_font = _map_font(preset.get("typography", "") if preset else "")

    # Heading on right
    right_left = half_width + Inches(0.6)
    right_width = half_width - Inches(1.2)
    top = Inches(2.0)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(right_left, top, right_width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = image_info["name"].title()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _hex_to_pptx_rgb(accent_color)
    p.font.name = heading_font
    p.alignment = PP_ALIGN.LEFT

    # Content text
    notes_text = summary_entry.get("prompt", "")
    if notes_text:
        content_top = Inches(3.8)
        content_height = Inches(2.5)
        txBox2 = slide.shapes.add_textbox(right_left, content_top, right_width, content_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = notes_text[:300]
        p2.font.size = Pt(16)
        p2.font.color.rgb = _hex_to_pptx_rgb(text_color)
        p2.font.name = body_font
        p2.alignment = PP_ALIGN.LEFT

    _add_logo(slide, prs, logo_path, preset)


_LAYOUT_BUILDERS = {
    "fullbleed": _add_content_slide_fullbleed,
    "standard": _add_content_slide_standard,
    "split": _add_content_slide_split,
}


def _add_content_slide(prs, image_info: dict, layout: str, preset, logo_path, summary_data):
    """Add one content slide using the specified layout."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    summary_entry = _get_summary_for_slide(summary_data, image_info["number"])
    builder = _LAYOUT_BUILDERS.get(layout, _add_content_slide_standard)
    builder(prs, slide, image_info, preset, logo_path, summary_entry)

    # Add slide notes with original prompt
    prompt = summary_entry.get("prompt", "")
    if prompt:
        slide.notes_slide.notes_text_frame.text = f"Prompt: {prompt}"

    return slide


def _add_closing_slide(prs, preset: dict | None, logo_path: str):
    """Create the closing 'Thank You' slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color, accent_color, text_color = _get_preset_colors(preset)
    _set_slide_bg(slide, bg_color)

    heading_font, body_font = _map_font(preset.get("typography", "") if preset else "")

    # Closing message
    closing_text = "Thank You"
    if preset and preset.get("mood"):
        closing_text = preset["mood"]

    left = Inches(1.5)
    top = Inches(2.5)
    width = prs.slide_width - Inches(3)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = closing_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = _hex_to_pptx_rgb(accent_color)
    p.font.name = heading_font
    p.alignment = PP_ALIGN.CENTER

    _add_logo(slide, prs, logo_path, preset)
    return slide


# ---------------------------------------------------------------------------
# Subcommands
# ---------------------------------------------------------------------------

def cmd_build(args):
    """Main build command -- assemble images into .pptx deck."""
    images_dir = Path(args.images).resolve()
    if not images_dir.is_dir():
        print(json.dumps({"error": True, "message": f"Images directory not found: {images_dir}"}))
        sys.exit(1)

    images = _load_slide_images(images_dir)
    if not images:
        print(json.dumps({"error": True, "message": f"No images found in {images_dir}"}))
        sys.exit(1)

    summary_data = _load_summary(images_dir)
    preset = _load_preset(args.preset) if args.preset else None
    layout = args.layout or DEFAULT_LAYOUT
    if layout not in LAYOUTS:
        print(json.dumps({"error": True, "message": f"Unknown layout: {layout}. Use: {', '.join(LAYOUTS)}"}))
        sys.exit(1)

    output_path = Path(args.output).resolve()

    # Build presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    title = args.title or "Presentation"
    subtitle = preset.get("tagline", "") if preset else ""
    logo_path = args.logo or ""

    _add_title_slide(prs, title, subtitle, preset, logo_path)

    for img_info in images:
        _add_content_slide(prs, img_info, layout, preset, logo_path, summary_data)

    _add_closing_slide(prs, preset, logo_path)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    total_slides = len(images) + 2  # content + title + closing
    result = {
        "output_path": str(output_path),
        "slides": total_slides,
        "layout": layout,
        "preset": args.preset or None,
        "title_slide": True,
        "closing_slide": True,
    }
    print(json.dumps(result, indent=2))


def cmd_estimate(args):
    """Count images and estimate output."""
    images_dir = Path(args.images).resolve()
    if not images_dir.is_dir():
        print(json.dumps({"error": True, "message": f"Images directory not found: {images_dir}"}))
        sys.exit(1)

    images = _load_slide_images(images_dir)
    count = len(images)
    total = count + 2  # title + closing

    result = {
        "images_found": count,
        "total_slides": total,
        "note": f"{count} content slides + title + closing",
    }
    print(json.dumps(result, indent=2))


def cmd_layouts(args):
    """Print available layout descriptions."""
    print(json.dumps({"layouts": LAYOUTS}, indent=2))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Banana Claude -- Slide Deck Builder",
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    # build
    p_build = subparsers.add_parser("build", help="Assemble slide images into .pptx")
    p_build.add_argument("--images", required=True, help="Directory containing slide images")
    p_build.add_argument("--output", required=True, help="Output .pptx file path")
    p_build.add_argument("--preset", default=None, help="Brand preset name from ~/.banana/presets/")
    p_build.add_argument("--title", default=None, help="Presentation title")
    p_build.add_argument("--layout", default=None, choices=list(LAYOUTS.keys()),
                         help="Slide layout mode (default: standard)")
    p_build.add_argument("--logo", default=None, help="Path to logo image file")
    p_build.set_defaults(func=cmd_build)

    # estimate
    p_est = subparsers.add_parser("estimate", help="Count images and estimate output")
    p_est.add_argument("--images", required=True, help="Directory containing slide images")
    p_est.set_defaults(func=cmd_estimate)

    # layouts
    p_lay = subparsers.add_parser("layouts", help="List available layout modes")
    p_lay.set_defaults(func=cmd_layouts)

    args = parser.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
