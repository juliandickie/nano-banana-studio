#!/usr/bin/env python3
"""Banana Claude -- Visual Brand Book Generator

Generate a complete visual brand book from a brand preset in three formats:
Markdown + images, PowerPoint (.pptx), and self-contained HTML.

Usage:
    brandbook.py generate --preset luxury-dark --output ~/brand-book/ --tier standard
    brandbook.py generate --preset luxury-dark --output ~/brand-book/ --tier quick --format html
"""

import argparse
import base64
import json
import os
import re
import sys
import time
import urllib.request
from datetime import datetime
from pathlib import Path

# python-pptx and Pillow are available
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Sibling import -- pantone_lookup lives next to this script
_SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(_SCRIPT_DIR))
from pantone_lookup import get_color_specs, hex_to_rgb

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_MODEL = "gemini-3.1-flash-image-preview"
API_BASE = "https://generativelanguage.googleapis.com/v1beta/models"
PRESETS_DIR = Path.home() / ".banana" / "presets"

PRICING = {
    "gemini-3.1-flash-image-preview": {"512": 0.020, "1K": 0.039, "2K": 0.078, "4K": 0.156},
}

# Image type -> (aspect_ratio, resolution)
IMAGE_SPECS = {
    "cover":                ("16:9", "4K"),
    "color-palette":        ("16:9", "4K"),
    "typography-specimen":  ("16:9", "4K"),
    "photo-product":        ("16:9", "4K"),
    "photo-editorial":      ("16:9", "4K"),
    "photo-portrait":       ("16:9", "4K"),
    "photo-environment":    ("16:9", "4K"),
    "motif-sample":         ("16:9", "4K"),
    "do-example-1":         ("16:9", "4K"),
    "do-example-2":         ("16:9", "4K"),
    "dont-example-1":       ("16:9", "4K"),
    "dont-example-2":       ("16:9", "4K"),
    "template-instagram":   ("4:5", "2K"),
    "template-linkedin":    ("16:9", "2K"),
    "template-slide-title": ("16:9", "4K"),
    "template-slide-content": ("16:9", "4K"),
}

# Tier definitions
TIERS = {
    "quick": [
        "cover", "color-palette", "typography-specimen",
        "photo-product", "photo-editorial",
    ],
    "standard": [
        "cover", "color-palette", "typography-specimen",
        "photo-product", "photo-editorial", "photo-portrait",
        "photo-environment", "motif-sample",
        "do-example-1", "do-example-2",
        "dont-example-1", "dont-example-2",
        "template-instagram", "template-linkedin",
        "template-slide-title", "template-slide-content",
    ],
    "comprehensive": [
        "cover", "color-palette", "typography-specimen",
        "photo-product", "photo-editorial", "photo-portrait",
        "photo-environment", "motif-sample",
        "do-example-1", "do-example-2", "do-example-3",
        "dont-example-1", "dont-example-2", "dont-example-3",
        "template-instagram", "template-linkedin",
        "template-slide-title", "template-slide-content",
        "photo-product-alt", "photo-editorial-alt",
        "photo-portrait-alt", "photo-environment-alt",
        "template-facebook", "template-twitter",
        "template-slide-divider", "template-slide-closing",
        "motif-sample-alt",
    ],
}

# Add comprehensive-only specs
for _extra in ["do-example-3", "dont-example-3", "photo-product-alt",
               "photo-editorial-alt", "photo-portrait-alt", "photo-environment-alt",
               "template-facebook", "template-twitter",
               "template-slide-divider", "template-slide-closing", "motif-sample-alt"]:
    if _extra.startswith("template-facebook"):
        IMAGE_SPECS[_extra] = ("4:5", "2K")
    elif _extra.startswith("template-twitter"):
        IMAGE_SPECS[_extra] = ("16:9", "2K")
    else:
        IMAGE_SPECS[_extra] = ("16:9", "4K")


# ---------------------------------------------------------------------------
# API key loading (same pattern as generate.py / slides.py)
# ---------------------------------------------------------------------------

def _load_api_key(cli_key=None):
    """Load Google AI API key from CLI, env, or config file."""
    key = cli_key or os.environ.get("GOOGLE_AI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not key:
        config_path = Path.home() / ".banana" / "config.json"
        if config_path.exists():
            try:
                with open(config_path) as f:
                    key = json.load(f).get("google_ai_api_key", "")
            except (json.JSONDecodeError, OSError):
                pass
    return key or None


# ---------------------------------------------------------------------------
# Image generation (same pattern as slides.py generate_single_image)
# ---------------------------------------------------------------------------

def generate_single_image(prompt, aspect_ratio, resolution, api_key,
                          model=DEFAULT_MODEL):
    """Call Gemini API to generate a single image. Returns (bytes, error)."""
    url = f"{API_BASE}/{model}:generateContent?key={api_key}"
    body = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": aspect_ratio,
                "imageSize": resolution,
            },
        },
    }
    data = json.dumps(body).encode("utf-8")

    max_retries = 3
    for attempt in range(max_retries):
        try:
            req = urllib.request.Request(
                url, data=data,
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=180) as resp:
                result = json.loads(resp.read().decode("utf-8"))

            candidates = result.get("candidates", [])
            if not candidates:
                reason = result.get("promptFeedback", {}).get("blockReason", "No candidates")
                return None, reason

            parts = candidates[0].get("content", {}).get("parts", [])
            for part in parts:
                if "inlineData" in part:
                    return base64.b64decode(part["inlineData"]["data"]), None

            finish_reason = candidates[0].get("finishReason", "UNKNOWN")
            return None, f"No image in response. finishReason: {finish_reason}"

        except urllib.error.HTTPError as e:
            error_body = e.read().decode("utf-8") if e.fp else ""
            if e.code == 429 and attempt < max_retries - 1:
                wait = 2 ** (attempt + 1)
                print(f"  Rate limited. Waiting {wait}s...", file=sys.stderr)
                time.sleep(wait)
                continue
            return None, f"HTTP {e.code}: {error_body[:200]}"
        except urllib.error.URLError as e:
            return None, f"Network error: {e.reason}"

    return None, "Max retries exceeded"


# ---------------------------------------------------------------------------
# Preset loading
# ---------------------------------------------------------------------------

def load_preset(name):
    """Load a preset JSON from ~/.banana/presets/NAME.json."""
    path = PRESETS_DIR / f"{name}.json"
    if not path.exists():
        # Try the bundled presets next to the script
        bundled = _SCRIPT_DIR.parent / "presets" / f"{name}.json"
        if bundled.exists():
            path = bundled
        else:
            print(json.dumps({"error": True, "message": f"Preset not found: {path} or {bundled}"}))
            sys.exit(1)
    with open(path) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# Prompt construction
# ---------------------------------------------------------------------------

def _keywords_str(preset):
    """Flatten prompt_keywords into a comma-separated string."""
    kw = preset.get("prompt_keywords", {})
    words = []
    for vals in kw.values():
        if isinstance(vals, list):
            words.extend(vals)
    return ", ".join(words[:8])


def _build_prompt(image_id, preset):
    """Build the generation prompt for a given image type."""
    brand = preset.get("name", "Brand")
    style = preset.get("style", "")
    mood = preset.get("mood", "")
    lighting = preset.get("lighting", "")
    typography = preset.get("typography", "")
    motifs = preset.get("visual_motifs", "")
    suffix = preset.get("prompt_suffix", "")
    colors_hex = preset.get("colors", [])
    do_list = preset.get("do_list", [])
    dont_list = preset.get("dont_list", [])
    kw = _keywords_str(preset)
    color_str = ", ".join(colors_hex[:5])

    base_id = re.sub(r'-(alt|\d+)$', '', image_id)

    prompts = {
        "cover": (
            f"A premium brand cover page for '{brand}'. {style}. {mood}. "
            f"{motifs} at low opacity in background. Dark background with brand "
            f"accent colors ({color_str}). {kw}. 16:9, 4K resolution. "
            f"NO text, NO logos. {suffix}"
        ),
        "color-palette": (
            f"A professional color palette visualization showing {len(colors_hex)} "
            f"color swatches arranged horizontally. Each swatch is a rounded "
            f"rectangle filled with a solid color: {color_str}. Clean dark "
            f"background, minimal design. 16:9, 4K resolution."
        ),
        "typography-specimen": (
            f"A typography specimen showing text hierarchy. Large headline in "
            f"{typography} reading 'HEADLINE', medium subheading reading "
            f"'Subheading Text', body paragraph reading 'Body text sample'. "
            f"{mood} mood. 16:9, 4K resolution. {suffix}"
        ),
        "photo-product": (
            f"A professional product photography shot. {style}. {lighting}. "
            f"Premium product centered on dark surface with brand color accents "
            f"({color_str}). {kw}. 16:9, 4K resolution. {suffix}"
        ),
        "photo-editorial": (
            f"An editorial photography shot. {style}. {lighting}. "
            f"Cinematic composition with negative space. {mood} mood. "
            f"{kw}. 16:9, 4K resolution. {suffix}"
        ),
        "photo-portrait": (
            f"A professional portrait photograph. {style}. {lighting}. "
            f"Subject lit with {lighting}. {mood} mood. Clean background. "
            f"{kw}. 16:9, 4K resolution. {suffix}"
        ),
        "photo-environment": (
            f"An environmental photograph showing an architectural interior. "
            f"{style}. {lighting}. Sophisticated space with brand color "
            f"accents ({color_str}). {mood} mood. {kw}. 16:9, 4K. {suffix}"
        ),
        "motif-sample": (
            f"A seamless pattern tile showing {motifs}. Brand colors: "
            f"{color_str}. Clean minimal design on dark background. "
            f"{kw}. 16:9, 4K resolution. {suffix}"
        ),
        "do-example": (
            f"A brand-correct example image following brand guidelines: "
            f"{'; '.join(do_list[:3])}. {style}. {lighting}. {mood} mood. "
            f"Brand colors ({color_str}). {kw}. 16:9, 4K. {suffix}"
        ),
        "dont-example": (
            f"A deliberately incorrect brand example that violates guidelines: "
            f"{'; '.join(dont_list[:3])}. Intentionally breaking the brand "
            f"rules to show what NOT to do. Cluttered, off-brand composition. "
            f"16:9, 4K resolution."
        ),
        "template-instagram": (
            f"A social media template for Instagram in portrait format. "
            f"{style}. Clean layout with large central area for imagery and "
            f"space for text at bottom. Brand colors ({color_str}). "
            f"{mood} mood. 4:5 ratio. NO text, NO logos. {suffix}"
        ),
        "template-linkedin": (
            f"A professional LinkedIn post template in landscape format. "
            f"{style}. Clean corporate layout with brand colors ({color_str}). "
            f"Large negative space for headline text. {mood} mood. "
            f"16:9 ratio. NO text, NO logos. {suffix}"
        ),
        "template-facebook": (
            f"A social media template for Facebook in portrait format. "
            f"{style}. Modern layout with brand colors ({color_str}). "
            f"Large visual area with text zone. {mood} mood. 4:5 ratio. "
            f"NO text, NO logos. {suffix}"
        ),
        "template-twitter": (
            f"A social media template for Twitter/X in landscape format. "
            f"{style}. Compact layout with brand colors ({color_str}). "
            f"{mood} mood. 16:9 ratio. NO text, NO logos. {suffix}"
        ),
        "template-slide-title": (
            f"A premium presentation title slide background. {style}. "
            f"{motifs} at 20% opacity. Dark background with brand accent "
            f"colors ({color_str}). Large clean negative space in center "
            f"for title text. 16:9, 4K. NO text, NO logos. {suffix}"
        ),
        "template-slide-content": (
            f"A premium presentation content slide background. {style}. "
            f"Subtle {motifs} at 15% opacity in corner. Mostly clean dark "
            f"space for content overlay. Brand accent ({color_str}). "
            f"16:9, 4K. NO text, NO logos. {suffix}"
        ),
        "template-slide-divider": (
            f"A premium presentation section divider slide background. "
            f"{style}. Bold use of brand accent color. Gradient sweep with "
            f"brand colors ({color_str}). Clean center for section title. "
            f"16:9, 4K. NO text, NO logos. {suffix}"
        ),
        "template-slide-closing": (
            f"A premium presentation closing slide background. {style}. "
            f"{motifs} at 25% opacity. Dark background with subtle brand "
            f"colors ({color_str}). {mood} mood. 16:9, 4K. "
            f"NO text, NO logos. {suffix}"
        ),
    }

    # Map compound IDs to base prompts
    if base_id.startswith("do-example"):
        return prompts["do-example"]
    if base_id.startswith("dont-example"):
        return prompts["dont-example"]

    # Handle -alt variants by reusing the base prompt with a variation note
    if image_id.endswith("-alt"):
        base_prompt = prompts.get(base_id, prompts.get("photo-product", ""))
        return base_prompt + " Alternative angle and composition."

    return prompts.get(image_id, prompts.get(base_id, f"{style}. {lighting}. {suffix}"))


# ---------------------------------------------------------------------------
# Generation plan + image generation
# ---------------------------------------------------------------------------

def build_generation_plan(preset, tier):
    """Return list of dicts describing each image to generate."""
    image_ids = TIERS.get(tier, TIERS["standard"])
    plan = []
    for img_id in image_ids:
        plan.append({
            "id": img_id,
            "prompt": _build_prompt(img_id, preset),
            "description": img_id.replace("-", " ").title(),
            "aspect_ratio": IMAGE_SPECS.get(img_id, ("16:9", "4K"))[0],
            "resolution": IMAGE_SPECS.get(img_id, ("16:9", "4K"))[1],
        })
    return plan


def generate_images(plan, preset, output_dir, api_key):
    """Generate all images in the plan. Returns {id: path_str} dict."""
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    total = len(plan)
    image_paths = {}
    succeeded = 0
    failed = 0

    for i, item in enumerate(plan):
        img_id = item["id"]
        ratio = item["aspect_ratio"]
        resolution = item["resolution"]
        filename = f"{img_id}.png"
        filepath = images_dir / filename

        print(f"  [{i+1}/{total}] Generating {item['description']}...", end=" ", flush=True)

        image_data, error = generate_single_image(
            prompt=item["prompt"],
            aspect_ratio=ratio,
            resolution=resolution,
            api_key=api_key,
        )

        if image_data:
            with open(filepath, "wb") as f:
                f.write(image_data)
            print(f"OK -> {filename}")
            image_paths[img_id] = str(filepath)
            succeeded += 1
        else:
            print(f"FAILED: {error}")
            failed += 1

        # Brief pause between generations
        if i < total - 1:
            time.sleep(1)

    return image_paths, succeeded, failed


# ---------------------------------------------------------------------------
# Color table
# ---------------------------------------------------------------------------

def build_color_table(colors):
    """Build full color spec table from a list of hex values."""
    table = []
    for hex_color in colors:
        if not hex_color.startswith("#"):
            hex_color = f"#{hex_color}"
        specs = get_color_specs(hex_color)
        table.append(specs)
    return table


# ---------------------------------------------------------------------------
# Output: Markdown
# ---------------------------------------------------------------------------

def output_markdown(preset, colors_table, image_paths, output_dir):
    """Write brand-book.md with images, color specs, and brand info."""
    brand = preset.get("name", "Brand")
    desc = preset.get("description", "")
    typography = preset.get("typography", "")
    style = preset.get("style", "")
    mood = preset.get("mood", "")
    motifs = preset.get("visual_motifs", "")
    do_list = preset.get("do_list", [])
    dont_list = preset.get("dont_list", [])

    lines = [
        f"# {brand} -- Visual Brand Book",
        "",
        f"> {desc}",
        "",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
    ]

    # Cover image
    if "cover" in image_paths:
        lines.append(f"![Cover](images/cover.png)")
        lines.append("")

    # Color Palette
    lines.append("## Color Palette")
    lines.append("")
    lines.append("| Swatch | Hex | RGB | CMYK | Pantone |")
    lines.append("|--------|-----|-----|------|---------|")
    for spec in colors_table:
        r, g, b = spec["rgb"]
        c, m, y, k = spec["cmyk"]
        lines.append(
            f"| {spec['hex']} | `{spec['hex']}` | "
            f"`rgb({r}, {g}, {b})` | "
            f"`cmyk({c}, {m}, {y}, {k})` | "
            f"{spec['pantone']} |"
        )
    lines.append("")

    if "color-palette" in image_paths:
        lines.append(f"![Color Palette](images/color-palette.png)")
        lines.append("")

    # Typography
    lines.append("## Typography")
    lines.append("")
    lines.append(f"{typography}")
    lines.append("")
    if "typography-specimen" in image_paths:
        lines.append(f"![Typography Specimen](images/typography-specimen.png)")
        lines.append("")

    # Photography Style
    lines.append("## Photography Style")
    lines.append("")
    lines.append(f"**Style:** {style}")
    lines.append("")
    lines.append(f"**Lighting:** {preset.get('lighting', '')}")
    lines.append("")
    lines.append(f"**Mood:** {mood}")
    lines.append("")
    for pid in ["photo-product", "photo-editorial", "photo-portrait", "photo-environment",
                "photo-product-alt", "photo-editorial-alt", "photo-portrait-alt", "photo-environment-alt"]:
        if pid in image_paths:
            label = pid.replace("-", " ").title()
            lines.append(f"### {label}")
            lines.append(f"![{label}](images/{pid}.png)")
            lines.append("")

    # Visual Motifs
    lines.append("## Visual Motifs")
    lines.append("")
    lines.append(f"{motifs}")
    lines.append("")
    for pid in ["motif-sample", "motif-sample-alt"]:
        if pid in image_paths:
            lines.append(f"![Motif Sample](images/{pid}.png)")
            lines.append("")

    # Do's and Don'ts
    lines.append("## Do's")
    lines.append("")
    for item in do_list:
        lines.append(f"- {item}")
    lines.append("")
    for pid in ["do-example-1", "do-example-2", "do-example-3"]:
        if pid in image_paths:
            lines.append(f"![Do Example](images/{pid}.png)")
            lines.append("")

    lines.append("## Don'ts")
    lines.append("")
    for item in dont_list:
        lines.append(f"- {item}")
    lines.append("")
    for pid in ["dont-example-1", "dont-example-2", "dont-example-3"]:
        if pid in image_paths:
            lines.append(f"![Don't Example](images/{pid}.png)")
            lines.append("")

    # Templates
    template_ids = [k for k in image_paths if k.startswith("template-")]
    if template_ids:
        lines.append("## Templates")
        lines.append("")
        for tid in template_ids:
            label = tid.replace("template-", "").replace("-", " ").title()
            lines.append(f"### {label}")
            lines.append(f"![{label}](images/{tid}.png)")
            lines.append("")

    md_path = output_dir / "brand-book.md"
    with open(md_path, "w") as f:
        f.write("\n".join(lines))
    return str(md_path)


# ---------------------------------------------------------------------------
# Output: PPTX
# ---------------------------------------------------------------------------

def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=(255, 255, 255), alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.alignment = alignment
    return txBox


def output_pptx(preset, colors_table, image_paths, output_dir, logo_path=None):
    """Create brand-book.pptx using python-pptx."""
    brand = preset.get("name", "Brand")
    colors_hex = preset.get("colors", [])
    typography = preset.get("typography", "")
    do_list = preset.get("do_list", [])
    dont_list = preset.get("dont_list", [])

    # Primary brand colors
    bg_rgb = hex_to_rgb(colors_hex[0]) if colors_hex else (0, 0, 0)
    accent_rgb = hex_to_rgb(colors_hex[1]) if len(colors_hex) > 1 else (255, 192, 0)
    text_rgb = hex_to_rgb(colors_hex[2]) if len(colors_hex) > 2 else (255, 255, 255)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank layout

    def _set_bg(slide, rgb):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)

    def _add_full_image(slide, img_path):
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                prs.slide_width, prs.slide_height,
            )

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(blank_layout)
    cover_img = image_paths.get("cover")
    if cover_img and Path(cover_img).exists():
        _add_full_image(slide, cover_img)
    else:
        _set_bg(slide, bg_rgb)
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                  brand.upper(), font_size=54, bold=True, color=text_rgb,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                  "Visual Brand Book", font_size=24, color=accent_rgb,
                  alignment=PP_ALIGN.CENTER)
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.5), Inches(6), Inches(1.5))

    # --- Slide 2: Color Palette ---
    slide = prs.slides.add_slide(blank_layout)
    _set_bg(slide, bg_rgb)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  "Color Palette", font_size=36, bold=True, color=text_rgb)
    swatch_w = Inches(2)
    swatch_h = Inches(2)
    start_x = Inches(0.5)
    for i, spec in enumerate(colors_table[:5]):
        x = start_x + i * Inches(2.4)
        r, g, b = spec["rgb"]
        shape = slide.shapes.add_shape(
            1, x, Inches(1.5), swatch_w, swatch_h,  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
        c, m, y, k = spec["cmyk"]
        info = f"{spec['hex']}\nrgb({r},{g},{b})\ncmyk({c},{m},{y},{k})\n{spec['pantone']}"
        _add_text_box(slide, x, Inches(3.7), swatch_w, Inches(2.5),
                      info, font_size=11, color=text_rgb)

    # --- Slide 3: Typography ---
    slide = prs.slides.add_slide(blank_layout)
    _set_bg(slide, bg_rgb)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  "Typography", font_size=36, bold=True, color=text_rgb)
    _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                  "HEADLINE", font_size=48, bold=True, color=accent_rgb)
    _add_text_box(slide, Inches(0.5), Inches(2.8), Inches(12), Inches(0.8),
                  "Subheading Text", font_size=28, color=text_rgb)
    _add_text_box(slide, Inches(0.5), Inches(3.8), Inches(10), Inches(1),
                  "Body text sample. " + typography, font_size=16, color=text_rgb)
    if "typography-specimen" in image_paths:
        img_path = image_paths["typography-specimen"]
        if Path(img_path).exists():
            slide.shapes.add_picture(
                img_path, Inches(0.5), Inches(5), Inches(5), Inches(2.2),
            )

    # --- Slides 4-7: Photography samples ---
    photo_ids = ["photo-product", "photo-editorial", "photo-portrait", "photo-environment",
                 "photo-product-alt", "photo-editorial-alt", "photo-portrait-alt", "photo-environment-alt"]
    for pid in photo_ids:
        if pid in image_paths:
            slide = prs.slides.add_slide(blank_layout)
            _set_bg(slide, bg_rgb)
            _add_full_image(slide, image_paths[pid])
            label = pid.replace("-", " ").title()
            _add_text_box(slide, Inches(0.5), Inches(6.5), Inches(5), Inches(0.6),
                          label, font_size=14, color=text_rgb)

    # --- Slide 8: Visual Motifs ---
    motif_ids = [k for k in image_paths if k.startswith("motif-")]
    if motif_ids:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg_rgb)
        _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Visual Motifs", font_size=36, bold=True, color=text_rgb)
        _add_text_box(slide, Inches(0.5), Inches(1.2), Inches(10), Inches(0.6),
                      preset.get("visual_motifs", ""), font_size=14, color=text_rgb)
        for j, mid in enumerate(motif_ids[:2]):
            if Path(image_paths[mid]).exists():
                slide.shapes.add_picture(
                    image_paths[mid], Inches(0.5), Inches(2 + j * 2.5),
                    Inches(8), Inches(2.2),
                )

    # --- Slides 9-10: Do's and Don'ts ---
    for label, items, example_prefix in [("Do's", do_list, "do-example"),
                                          ("Don'ts", dont_list, "dont-example")]:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg_rgb)
        _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      label, font_size=36, bold=True, color=text_rgb)
        rules_text = "\n".join(f"  {item}" for item in items)
        _add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5),
                      rules_text, font_size=14, color=text_rgb)
        # Place example images on right side
        examples = [k for k in image_paths if k.startswith(example_prefix)]
        for j, eid in enumerate(examples[:2]):
            if Path(image_paths[eid]).exists():
                slide.shapes.add_picture(
                    image_paths[eid], Inches(6.5), Inches(1.3 + j * 2.8),
                    Inches(6), Inches(2.5),
                )

    # --- Slides 11+: Templates ---
    template_ids = sorted(k for k in image_paths if k.startswith("template-"))
    for tid in template_ids:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg_rgb)
        _add_full_image(slide, image_paths[tid])
        label = tid.replace("template-", "").replace("-", " ").title()
        _add_text_box(slide, Inches(0.5), Inches(6.5), Inches(5), Inches(0.6),
                      f"Template: {label}", font_size=14, color=text_rgb)

    pptx_path = output_dir / "brand-book.pptx"
    prs.save(str(pptx_path))
    return str(pptx_path)


# ---------------------------------------------------------------------------
# Output: HTML
# ---------------------------------------------------------------------------

def _img_to_data_uri(img_path):
    """Read an image file and return a base64 data URI."""
    if not img_path or not Path(img_path).exists():
        return ""
    with open(img_path, "rb") as f:
        data = base64.b64encode(f.read()).decode("ascii")
    return f"data:image/png;base64,{data}"


def output_html(preset, colors_table, image_paths, output_dir, logo_path=None):
    """Create a self-contained brand-book.html with base64 images."""
    brand = preset.get("name", "Brand")
    desc = preset.get("description", "")
    typography = preset.get("typography", "")
    style = preset.get("style", "")
    mood = preset.get("mood", "")
    motifs = preset.get("visual_motifs", "")
    do_list = preset.get("do_list", [])
    dont_list = preset.get("dont_list", [])
    colors_hex = preset.get("colors", [])

    bg = colors_hex[0] if colors_hex else "#000000"
    accent = colors_hex[1] if len(colors_hex) > 1 else "#FFC000"
    text_color = colors_hex[2] if len(colors_hex) > 2 else "#FFFFFF"

    # Build color rows
    color_rows = ""
    for spec in colors_table:
        r, g, b = spec["rgb"]
        c, m, y, k = spec["cmyk"]
        color_rows += f"""
        <tr>
          <td><div class="swatch" style="background:{spec['hex']}"></div></td>
          <td><code>{spec['hex']}</code></td>
          <td><code>rgb({r}, {g}, {b})</code></td>
          <td><code>cmyk({c}, {m}, {y}, {k})</code></td>
          <td>{spec['pantone']}</td>
        </tr>"""

    # Build image sections
    def _img_tag(img_id, alt=""):
        uri = _img_to_data_uri(image_paths.get(img_id))
        if not uri:
            return ""
        return f'<img src="{uri}" alt="{alt}" class="brand-img" />'

    photo_html = ""
    for pid in ["photo-product", "photo-editorial", "photo-portrait", "photo-environment",
                "photo-product-alt", "photo-editorial-alt", "photo-portrait-alt", "photo-environment-alt"]:
        tag = _img_tag(pid, pid)
        if tag:
            label = pid.replace("-", " ").title()
            photo_html += f'<div class="photo-card"><h4>{label}</h4>{tag}</div>\n'

    motif_html = ""
    for mid in [k for k in image_paths if k.startswith("motif-")]:
        tag = _img_tag(mid, "Motif Sample")
        if tag:
            motif_html += tag + "\n"

    do_items = "\n".join(f"<li>{item}</li>" for item in do_list)
    dont_items = "\n".join(f"<li>{item}</li>" for item in dont_list)

    do_images = ""
    for pid in ["do-example-1", "do-example-2", "do-example-3"]:
        tag = _img_tag(pid, "Do Example")
        if tag:
            do_images += tag + "\n"

    dont_images = ""
    for pid in ["dont-example-1", "dont-example-2", "dont-example-3"]:
        tag = _img_tag(pid, "Don't Example")
        if tag:
            dont_images += tag + "\n"

    template_html = ""
    for tid in sorted(k for k in image_paths if k.startswith("template-")):
        tag = _img_tag(tid, tid)
        if tag:
            label = tid.replace("template-", "").replace("-", " ").title()
            template_html += f'<div class="template-card"><h4>{label}</h4>{tag}</div>\n'

    logo_img = ""
    if logo_path and Path(logo_path).exists():
        with open(logo_path, "rb") as f:
            logo_b64 = base64.b64encode(f.read()).decode("ascii")
        ext = Path(logo_path).suffix.lstrip(".")
        if ext in ("svg",):
            mime = "image/svg+xml"
        else:
            mime = f"image/{ext}" if ext else "image/png"
        logo_img = f'<img src="data:{mime};base64,{logo_b64}" alt="Logo" class="logo" />'

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8" />
<meta name="viewport" content="width=device-width, initial-scale=1.0" />
<title>{brand} -- Visual Brand Book</title>
<style>
  :root {{
    --bg: {bg};
    --accent: {accent};
    --text: {text_color};
  }}
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{
    background: var(--bg); color: var(--text);
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Helvetica, Arial, sans-serif;
    line-height: 1.6; padding: 0;
  }}
  .container {{ max-width: 1200px; margin: 0 auto; padding: 2rem; }}
  h1 {{ font-size: 3rem; font-weight: 700; letter-spacing: 0.05em; margin-bottom: 0.5rem; }}
  h2 {{
    font-size: 1.8rem; color: var(--accent); margin: 3rem 0 1rem;
    border-bottom: 2px solid var(--accent); padding-bottom: 0.5rem;
  }}
  h3 {{ font-size: 1.3rem; margin: 1.5rem 0 0.5rem; }}
  h4 {{ font-size: 1rem; margin-bottom: 0.5rem; color: var(--accent); }}
  p, li {{ font-size: 1rem; margin-bottom: 0.5rem; }}
  code {{ background: rgba(255,255,255,0.1); padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }}
  .brand-img {{ width: 100%; max-width: 900px; border-radius: 8px; margin: 1rem 0; display: block; }}
  .cover-section {{
    text-align: center; padding: 4rem 2rem;
    position: relative; overflow: hidden;
  }}
  .cover-section img.brand-img {{ max-width: 100%; border-radius: 0; }}
  .cover-section .overlay {{
    position: relative; z-index: 1; padding: 2rem;
  }}
  .subtitle {{ font-size: 1.3rem; color: var(--accent); margin-top: 0.5rem; }}
  .logo {{ max-height: 60px; margin-bottom: 1rem; }}
  table {{ width: 100%; border-collapse: collapse; margin: 1rem 0; }}
  th, td {{ padding: 10px 14px; text-align: left; border-bottom: 1px solid rgba(255,255,255,0.15); }}
  th {{ color: var(--accent); font-weight: 600; }}
  .swatch {{ width: 40px; height: 40px; border-radius: 6px; border: 1px solid rgba(255,255,255,0.2); }}
  .photo-grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(400px, 1fr)); gap: 1.5rem; }}
  .photo-card, .template-card {{ background: rgba(255,255,255,0.03); border-radius: 8px; padding: 1rem; }}
  .dos-donts {{ display: grid; grid-template-columns: 1fr 1fr; gap: 2rem; margin-top: 1rem; }}
  .dos-donts ul {{ padding-left: 1.2rem; }}
  .examples-row {{ display: flex; gap: 1rem; flex-wrap: wrap; margin-top: 1rem; }}
  .examples-row img {{ max-width: 420px; }}
  .template-grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(350px, 1fr)); gap: 1.5rem; }}
  .generated {{ text-align: center; font-size: 0.85rem; opacity: 0.5; margin-top: 3rem; }}
  @media print {{
    body {{ background: white; color: #222; }}
    h2 {{ color: {accent}; border-color: {accent}; }}
    .swatch {{ border-color: #ccc; }}
    .brand-img {{ max-width: 100%; page-break-inside: avoid; }}
  }}
</style>
</head>
<body>
<div class="container">

  <section class="cover-section">
    {logo_img}
    {_img_tag("cover", "Cover")}
    <div class="overlay">
      <h1>{brand}</h1>
      <p class="subtitle">Visual Brand Book</p>
      <p>{desc}</p>
    </div>
  </section>

  <h2>Color Palette</h2>
  {_img_tag("color-palette", "Color Palette")}
  <table>
    <thead><tr><th>Swatch</th><th>Hex</th><th>RGB</th><th>CMYK</th><th>Pantone</th></tr></thead>
    <tbody>{color_rows}</tbody>
  </table>

  <h2>Typography</h2>
  <p>{typography}</p>
  {_img_tag("typography-specimen", "Typography Specimen")}

  <h2>Photography Style</h2>
  <p><strong>Style:</strong> {style}</p>
  <p><strong>Mood:</strong> {mood}</p>
  <div class="photo-grid">{photo_html}</div>

  <h2>Visual Motifs</h2>
  <p>{motifs}</p>
  {motif_html}

  <h2>Do's and Don'ts</h2>
  <div class="dos-donts">
    <div>
      <h3>Do's</h3>
      <ul>{do_items}</ul>
      <div class="examples-row">{do_images}</div>
    </div>
    <div>
      <h3>Don'ts</h3>
      <ul>{dont_items}</ul>
      <div class="examples-row">{dont_images}</div>
    </div>
  </div>

  <h2>Templates</h2>
  <div class="template-grid">{template_html}</div>

  <p class="generated">Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} by Banana Claude</p>
</div>
</body>
</html>"""

    html_path = output_dir / "brand-book.html"
    with open(html_path, "w") as f:
        f.write(html)
    return str(html_path)


# ---------------------------------------------------------------------------
# Main command: generate
# ---------------------------------------------------------------------------

def cmd_generate(args):
    """Generate a visual brand book from a preset."""
    preset = load_preset(args.preset)

    api_key = _load_api_key(args.api_key)
    if not api_key:
        print(json.dumps({"error": True, "message": "No API key. Run /create-image setup, set GOOGLE_AI_API_KEY env, or pass --api-key"}))
        sys.exit(1)

    tier = args.tier or "standard"
    fmt = args.format or "all"
    output_dir = Path(args.output).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    logo_path = args.logo if args.logo else None

    # Build plan
    plan = build_generation_plan(preset, tier)
    total_images = len(plan)
    cost_per = PRICING[DEFAULT_MODEL].get("4K", 0.156)
    est_cost = sum(
        PRICING[DEFAULT_MODEL].get(item["resolution"], 0.156) for item in plan
    )

    print(f"Brand Book: {preset.get('name', args.preset)}")
    print(f"  Tier: {tier} ({total_images} images)")
    print(f"  Format: {fmt}")
    print(f"  Output: {output_dir}")
    print(f"  Estimated cost: ${est_cost:.2f}")
    print()

    # Generate images
    image_paths, succeeded, failed = generate_images(plan, preset, output_dir, api_key)

    print()
    print(f"Images: {succeeded}/{total_images} generated ({failed} failed)")

    # Build color table
    colors_table = build_color_table(preset.get("colors", []))

    # Output formats
    outputs = {}
    if fmt in ("all", "markdown"):
        md_path = output_markdown(preset, colors_table, image_paths, output_dir)
        outputs["markdown"] = md_path
        print(f"  Markdown: {md_path}")

    if fmt in ("all", "pptx"):
        pptx_path = output_pptx(preset, colors_table, image_paths, output_dir, logo_path)
        outputs["pptx"] = pptx_path
        print(f"  PPTX: {pptx_path}")

    if fmt in ("all", "html"):
        html_path = output_html(preset, colors_table, image_paths, output_dir, logo_path)
        outputs["html"] = html_path
        print(f"  HTML: {html_path}")

    # Write generation summary
    actual_cost = sum(
        PRICING[DEFAULT_MODEL].get(
            IMAGE_SPECS.get(img_id, ("16:9", "4K"))[1], 0.156
        )
        for img_id in image_paths
    )
    summary = {
        "preset": args.preset,
        "tier": tier,
        "format": fmt,
        "total_images": total_images,
        "succeeded": succeeded,
        "failed": failed,
        "estimated_cost": round(actual_cost, 3),
        "images": list(image_paths.keys()),
        "outputs": outputs,
        "generated_at": datetime.now().isoformat(),
    }
    summary_path = output_dir / "generation-summary.json"
    with open(summary_path, "w") as f:
        json.dump(summary, f, indent=2)

    print()
    print(f"Done! Summary: {summary_path}")
    print(json.dumps({
        "success": True,
        "preset": args.preset,
        "tier": tier,
        "total": total_images,
        "succeeded": succeeded,
        "failed": failed,
        "output_dir": str(output_dir),
        "estimated_cost": round(actual_cost, 3),
        "outputs": outputs,
    }))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Banana Claude -- Visual Brand Book Generator",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Examples:
  %(prog)s generate --preset luxury-dark --output ~/brand-book/ --tier standard
  %(prog)s generate --preset luxury-dark --output ~/brand-book/ --tier quick --format html
  %(prog)s generate --preset tech-saas --output ~/brand-book/ --tier comprehensive --format pptx
""",
    )
    sub = parser.add_subparsers(dest="command", required=True)

    p_gen = sub.add_parser("generate", help="Generate a visual brand book from a preset")
    p_gen.add_argument("--preset", required=True, help="Brand preset name (e.g. luxury-dark)")
    p_gen.add_argument("--output", required=True, help="Output directory")
    p_gen.add_argument("--tier", choices=["quick", "standard", "comprehensive"],
                       default="standard", help="Image tier (default: standard)")
    p_gen.add_argument("--format", choices=["all", "markdown", "pptx", "html"],
                       default="all", help="Output format (default: all)")
    p_gen.add_argument("--logo", default=None, help="Path to logo image file to embed")
    p_gen.add_argument("--api-key", default=None, help="Google AI API key")

    args = parser.parse_args()
    if args.command == "generate":
        cmd_generate(args)


if __name__ == "__main__":
    main()
